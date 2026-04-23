"""
文档格式转换器模块 (parsers.py)  —  LlamaIndex 重构版

核心变化：
- 原始手写解析器 → LlamaIndex 官方 Reader（llama-index-readers-file）
- 各 Reader 返回 List[LlamaDocument]，封装到业务 Document 对象
- 保持与旧版完全相同的公开接口：
    get_parser_for_file(file_path, config) -> BaseParser
    parse_document(file_path, config)      -> Document

支持格式及对应 Reader：
┌──────────┬─────────────────────────────────────────────────┐
│ 格式      │ LlamaIndex Reader                               │
├──────────┼─────────────────────────────────────────────────┤
│ PDF       │ PyMuPDFReader（优先）/ PDFReader（降级）         │
│ DOCX/DOC  │ DocxReader                                      │
│ XLSX/XLS  │ PandasExcelReader + 行转句（自定义补充）         │
│ CSV       │ PandasCSVReader                                 │
│ MD        │ FlatReader（原生 Markdown 保留结构）             │
│ TXT       │ SimpleDirectoryReader / FlatReader              │
└──────────┴─────────────────────────────────────────────────┘

所有 Reader 输出最终封装为 Document（models.py），
markdown_content 字段存放 Reader 提取的文本。
"""

import logging
import os
import re
from abc import ABC, abstractmethod
from typing import Any, Dict, List, Optional, Tuple

from llama_index.core.schema import Document as LlamaDocument

from .models import ChunkMetadata, Document

logger = logging.getLogger("oa_agent.parsers")


# ─────────────────────────────────────────────────────────────────────────────
# 基类
# ─────────────────────────────────────────────────────────────────────────────

class BaseParser(ABC):
    """
    解析器基类（接口与旧版完全一致）。

    子类只需实现 parse()，返回 (markdown_text, [ChunkMetadata, ...])。
    """

    def __init__(self, config: Optional[Dict[str, Any]] = None):
        self.config = config or {}

    @abstractmethod
    def parse(self, file_path: str) -> Tuple[str, List[ChunkMetadata]]:
        """
        解析文件。

        返回：
        - markdown_content: 提取的文本内容（Markdown 格式或纯文本）
        - chunk_metadata_list: 每个页面/章节的元数据
        """

    # ── 工具方法（子类可复用） ───────────────────────────────────────────────

    def clean_markdown(self, content: str) -> str:
        """清洗 Markdown 文本（去页眉页脚、水印、多余空行）"""
        content = re.sub(r"\n{3,}", "\n\n", content)
        content = content.replace("\t", "    ")
        for pattern in [
            r"第\s*\d+\s*页\s*[/／]\s*共\s*\d+\s*页",
            r"Page\s+\d+\s+of\s+\d+",
            r"-\s*\d+\s*-",
        ]:
            content = re.sub(pattern, "", content, flags=re.IGNORECASE)
        for kw in ["CONFIDENTIAL", "内部文件", "草稿", "DRAFT"]:
            content = re.sub(
                rf"^{kw}.*$", "", content, flags=re.IGNORECASE | re.MULTILINE
            )
        return content.strip()

    def extract_header_path(self, markdown_content: str) -> List[Tuple[str, str]]:
        """从 Markdown 提取标题层级路径列表 [(path, content), ...]"""
        lines = markdown_content.split("\n")
        path_parts: List[str] = []
        blocks: List[Tuple[str, str]] = []
        current_block: List[str] = []

        for line in lines:
            m = re.match(r"^(#{1,6})\s+(.+)$", line)
            if m:
                if current_block:
                    blocks.append((" -> ".join(path_parts), "\n".join(current_block)))
                    current_block = []
                level = len(m.group(1))
                title = m.group(2).strip()
                path_parts = path_parts[: level - 1] if level <= len(path_parts) else path_parts
                path_parts.append(title)
            else:
                current_block.append(line)

        if current_block:
            blocks.append((" -> ".join(path_parts), "\n".join(current_block)))

        return blocks

    def _llama_docs_to_text(self, llama_docs: List[LlamaDocument]) -> str:
        """将多个 LlamaDocument 拼接为单一文本"""
        return "\n\n".join(d.get_content() for d in llama_docs if d.get_content())


# ─────────────────────────────────────────────────────────────────────────────
# PDF 解析器
# ─────────────────────────────────────────────────────────────────────────────

class PDFParser(BaseParser):
    """
    PDF 解析器

    优先使用 PyMuPDFReader（精确提取页面文本 + 页码元数据），
    降级到 PDFReader（基于 pypdf）。
    """

    def parse(self, file_path: str) -> Tuple[str, List[ChunkMetadata]]:
        llama_docs = self._load_with_pymupdf(file_path)
        if not llama_docs:
            llama_docs = self._load_with_pdfreader(file_path)
        if not llama_docs:
            logger.error(f"PDF 解析失败，无可用 Reader：{file_path}")
            return "", []

        markdown_parts: List[str] = []
        metadata_list: List[ChunkMetadata] = []

        for doc in llama_docs:
            page_text = doc.get_content()
            if not page_text.strip():
                continue
            page_num = doc.metadata.get("page_label") or doc.metadata.get("page", "?")
            markdown_parts.append(f"## 第 {page_num} 页\n\n{page_text}")
            metadata_list.append(
                ChunkMetadata(
                    page_number=int(page_num) if str(page_num).isdigit() else None,
                    source_file=os.path.basename(file_path),
                )
            )

        full_markdown = "\n\n".join(markdown_parts)
        return self.clean_markdown(full_markdown), metadata_list

    def _load_with_pymupdf(self, file_path: str) -> List[LlamaDocument]:
        try:
            from llama_index.readers.file import PyMuPDFReader
            reader = PyMuPDFReader()
            return reader.load(file_path)
        except ImportError:
            logger.warning("llama-index-readers-file[pymupdf] 未安装，尝试降级")
        except Exception as e:
            logger.warning(f"PyMuPDFReader 失败: {e}")
        return []

    def _load_with_pdfreader(self, file_path: str) -> List[LlamaDocument]:
        try:
            from llama_index.readers.file import PDFReader
            reader = PDFReader()
            return reader.load_data(file_path)
        except ImportError:
            logger.warning("PDFReader 依赖缺失，请安装: pip install pypdf")
        except Exception as e:
            logger.warning(f"PDFReader 失败: {e}")
        return []


# ─────────────────────────────────────────────────────────────────────────────
# Word 解析器
# ─────────────────────────────────────────────────────────────────────────────

class WordParser(BaseParser):
    """
    Word 解析器

    使用 DocxReader 提取文档文本，保持标题层级结构。
    """

    def parse(self, file_path: str) -> Tuple[str, List[ChunkMetadata]]:
        llama_docs = self._load_with_docx_reader(file_path)
        if not llama_docs:
            return "", []

        full_text = self._llama_docs_to_text(llama_docs)
        full_text = self.clean_markdown(full_text)

        metadata_list: List[ChunkMetadata] = []
        blocks = self.extract_header_path(full_text)
        for path, _ in blocks:
            metadata_list.append(
                ChunkMetadata(
                    header_path=path,
                    source_file=os.path.basename(file_path),
                )
            )
        if not metadata_list:
            metadata_list = [ChunkMetadata(source_file=os.path.basename(file_path))]

        return full_text, metadata_list

    def _load_with_docx_reader(self, file_path: str) -> List[LlamaDocument]:
        try:
            from llama_index.readers.file import DocxReader
            reader = DocxReader()
            return reader.load_data(file_path)
        except ImportError:
            logger.error(
                "DocxReader 依赖缺失，请安装: pip install llama-index-readers-file docx2txt"
            )
        except Exception as e:
            logger.error(f"DocxReader 失败: {e}")
        return []


# ─────────────────────────────────────────────────────────────────────────────
# Excel 解析器
# ─────────────────────────────────────────────────────────────────────────────

class ExcelParser(BaseParser):
    """
    Excel 解析器

    流程：
    1. PandasExcelReader 读取所有 Sheet → LlamaDocument
    2. 同时用 pandas 做行转句（保持与旧版一致的自然语言描述）
    3. 可选：原始表格存入 PostgreSQL（需配置 save_to_rdb=True）
    """

    def __init__(self, config: Optional[Dict[str, Any]] = None):
        super().__init__(config)
        self.save_to_rdb: bool = self.config.get("save_to_rdb", True)

    def parse(self, file_path: str) -> Tuple[str, List[ChunkMetadata]]:
        try:
            import pandas as pd
        except ImportError:
            logger.error("pandas 未安装：pip install pandas openpyxl")
            return "", []

        ext = os.path.splitext(file_path)[1].lower()
        all_sheets_md: List[str] = []
        metadata_list: List[ChunkMetadata] = []

        try:
            if ext == ".csv":
                df_map = {"Sheet1": pd.read_csv(file_path)}
            else:
                excel_file = pd.ExcelFile(file_path)
                df_map = {
                    sheet: pd.read_excel(excel_file, sheet_name=sheet)
                    for sheet in excel_file.sheet_names
                }
        except Exception as e:
            logger.error(f"Excel 读取失败: {e}")
            return "", []

        for sheet_name, df in df_map.items():
            table_metadata = None
            if self.save_to_rdb:
                try:
                    from .excel_storage import get_excel_storage
                    doc_id = os.path.basename(file_path)
                    storage = get_excel_storage()
                    table_metadata = storage.save_table(doc_id, df, sheet_name)
                    logger.info(f"Excel 已存入 PostgreSQL: {table_metadata.table_name}")
                except Exception as e:
                    logger.warning(f"PostgreSQL 存储跳过: {e}")

            sentences = self._rows_to_sentences(df, sheet_name)
            sheet_md = [f"## {sheet_name}\n"] + sentences
            all_sheets_md.append("\n\n".join(sheet_md))

            extra: Dict[str, Any] = {
                "sheet_name": sheet_name,
                "row_count": len(df),
                "column_count": len(df.columns),
                "columns": list(df.columns),
            }
            if table_metadata:
                extra["table_name"] = table_metadata.table_name
                extra["has_relational_data"] = True
                extra["stat_columns"] = table_metadata.stat_columns

            metadata_list.append(
                ChunkMetadata(source_file=os.path.basename(file_path), extra=extra)
            )

        full_markdown = "\n\n".join(all_sheets_md)
        return self.clean_markdown(full_markdown), metadata_list

    def _rows_to_sentences(self, df: "pd.DataFrame", sheet_name: str) -> List[str]:  # type: ignore[name-defined]
        """DataFrame 每行转自然语言描述"""
        try:
            import pandas as pd
        except ImportError:
            return []
        sentences: List[str] = []
        columns = df.columns.tolist()
        for idx, row in df.iterrows():
            parts = []
            for col, val in zip(columns, row):
                if pd.notna(val) and str(val).strip():
                    parts.append(f"{str(col).strip()}为{str(val).strip()}")
            if parts:
                sentences.append(
                    f"在《{sheet_name}》中，第{idx + 1}行：{'，'.join(parts)}。"
                )
        return sentences


# ─────────────────────────────────────────────────────────────────────────────
# Markdown 解析器
# ─────────────────────────────────────────────────────────────────────────────

class MarkdownParser(BaseParser):
    """
    Markdown 解析器

    使用 LlamaIndex FlatReader 加载文件，保持原始 Markdown 结构。
    """

    def parse(self, file_path: str) -> Tuple[str, List[ChunkMetadata]]:
        content = self._read_file(file_path)
        if not content:
            return "", []

        content = self.clean_markdown(content)
        metadata_list: List[ChunkMetadata] = []
        for path, _ in self.extract_header_path(content):
            metadata_list.append(
                ChunkMetadata(
                    header_path=path,
                    source_file=os.path.basename(file_path),
                )
            )
        if not metadata_list:
            metadata_list = [ChunkMetadata(source_file=os.path.basename(file_path))]

        return content, metadata_list

    def _read_file(self, file_path: str) -> str:
        try:
            from llama_index.readers.file import FlatReader
            from pathlib import Path
            docs = FlatReader().load_data(Path(file_path))
            if docs:
                return self._llama_docs_to_text(docs)
        except ImportError:
            logger.warning("FlatReader 不可用，回退到直接读取")
        except Exception as e:
            logger.warning(f"FlatReader 失败: {e}")

        for enc in ("utf-8", "gbk"):
            try:
                with open(file_path, "r", encoding=enc) as f:
                    return f.read()
            except UnicodeDecodeError:
                continue
            except Exception as e:
                logger.error(f"读取 Markdown 文件失败: {e}")
                break
        return ""


# ─────────────────────────────────────────────────────────────────────────────
# 文本解析器
# ─────────────────────────────────────────────────────────────────────────────

class TextParser(BaseParser):
    """
    纯文本解析器

    使用 FlatReader 加载 .txt 文件。
    精细切分由 SemanticChunker（NodeParser）处理。
    """

    def parse(self, file_path: str) -> Tuple[str, List[ChunkMetadata]]:
        content = self._read_file(file_path)
        if not content:
            return "", []

        return content, [ChunkMetadata(source_file=os.path.basename(file_path))]

    def _read_file(self, file_path: str) -> str:
        try:
            from llama_index.readers.file import FlatReader
            from pathlib import Path
            docs = FlatReader().load_data(Path(file_path))
            if docs:
                return "\n\n".join(d.get_content() for d in docs)
        except ImportError:
            pass
        except Exception as e:
            logger.warning(f"FlatReader 失败: {e}")

        for enc in ("utf-8", "gbk"):
            try:
                with open(file_path, "r", encoding=enc) as f:
                    return f.read()
            except UnicodeDecodeError:
                continue
            except Exception as e:
                logger.error(f"读取文本文件失败: {e}")
                break
        return ""


# ─────────────────────────────────────────────────────────────────────────────
# PowerPoint 解析器
# ─────────────────────────────────────────────────────────────────────────────

class PPTXParser(BaseParser):
    """
    PowerPoint 解析器

    使用 python-pptx 提取幻灯片文本，保持标题层级结构。
    支持提取每张幻灯片的标题和内容。
    """

    def parse(self, file_path: str) -> Tuple[str, List[ChunkMetadata]]:
        content = self._read_pptx(file_path)
        if not content:
            return "", []

        content = self.clean_markdown(content)
        metadata_list: List[ChunkMetadata] = []
        for path, _ in self.extract_header_path(content):
            metadata_list.append(
                ChunkMetadata(
                    header_path=path,
                    source_file=os.path.basename(file_path),
                )
            )
        if not metadata_list:
            metadata_list = [ChunkMetadata(source_file=os.path.basename(file_path))]

        return content, metadata_list

    def _read_pptx(self, file_path: str) -> str:
        """使用 python-pptx 读取 PowerPoint 文件"""
        try:
            from pptx import Presentation
        except ImportError:
            logger.error("python-pptx 未安装：pip install python-pptx")
            return ""
        except Exception as e:
            logger.error(f"python-pptx 导入失败: {e}")
            return ""

        try:
            prs = Presentation(file_path)
            slides_content = []

            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_texts = []

                # 提取标题
                if slide.shapes.title:
                    title_text = slide.shapes.title.text.strip()
                    if title_text:
                        slide_texts.append(f"## 幻灯片 {slide_num}: {title_text}")
                    else:
                        slide_texts.append(f"## 幻灯片 {slide_num}")
                else:
                    slide_texts.append(f"## 幻灯片 {slide_num}")

                # 提取内容（排除标题形状，避免重复）
                title_shape = slide.shapes.title
                for shape in slide.shapes:
                    if shape == title_shape:
                        continue  # 跳过标题，避免重复遍历
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                slide_texts.append(f"- {text}")

                if len(slide_texts) > 1:  # 如果有除标题外的内容
                    slides_content.append("\n".join(slide_texts))

            return "\n\n".join(slides_content)

        except Exception as e:
            logger.error(f"PPTX 解析失败: {e}")
            return ""


# ─────────────────────────────────────────────────────────────────────────────
# 工厂函数
# ─────────────────────────────────────────────────────────────────────────────

_PARSER_MAP: Dict[str, type] = {
    "pdf": PDFParser,
    "docx": WordParser,
    "doc": WordParser,
    "xlsx": ExcelParser,
    "xls": ExcelParser,
    "csv": ExcelParser,
    "pptx": PPTXParser,
    "md": MarkdownParser,
    "markdown": MarkdownParser,
    "txt": TextParser,
}


def get_parser_for_file(
    file_path: str,
    config: Optional[Dict[str, Any]] = None,
) -> BaseParser:
    """
    工厂函数：根据文件扩展名获取对应解析器实例。

    未知扩展名使用 TextParser 兜底。
    """
    ext = os.path.splitext(file_path)[1].lower().lstrip(".")
    parser_cls = _PARSER_MAP.get(ext, TextParser)
    return parser_cls(config)


def parse_document(
    file_path: str,
    config: Optional[Dict[str, Any]] = None,
) -> Document:
    """
    便捷函数：将文件解析为 Document 对象。

    流程：
    1. 选取对应解析器
    2. 执行解析，得到 markdown_content
    3. 封装为 Document 对象（同时生成 llama_doc）
    """
    parser = get_parser_for_file(file_path, config)
    markdown_content, _ = parser.parse(file_path)

    doc = Document.create(file_path=file_path)
    doc.markdown_content = markdown_content

    llama_doc = LlamaDocument(
        doc_id=doc.doc_id,
        text=markdown_content,
        metadata={
            "doc_id": doc.doc_id,
            "title": doc.title,
            "format": doc.format.value,
            "file_path": doc.file_path,
            "content_hash": doc.content_hash,
            "file_size": doc.file_size,
        },
    )
    doc.llama_doc = llama_doc

    return doc
